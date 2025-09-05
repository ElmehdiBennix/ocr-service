
import csv
import magic
import pytesseract
import pdfplumber
import docx
import openpyxl
import pptx
import ebooklib
from ebooklib import epub
from striprtf.striprtf import rtf_to_text
from bs4 import BeautifulSoup
import re
from markdown_it import MarkdownIt
import os

from odf import text, teletype
from odf.opendocument import load

# custom Exceptions
class UnsupportedFileType(Exception):
    """Raised when the file type is not supported."""
    pass

class ParsingError(Exception):
    """Raised when there is an error parsing the file."""
    pass

def clean_text(text: str) -> str:
    """
    Cleans the extracted text by:
    - Removing extra whitespace and newlines.
    - Normalizing line breaks.
    """
    # replace multiple newlines with a single one
    text = re.sub(r'\n+', '\n', text)
    # strip leading/trailing whitespace from each line
    text = "\n".join([line.strip() for line in text.split('\n')])
    # strip leading/trailing whitespace from the whole text
    return text.strip()

def parse_pdf(file_path: str) -> str:
    """Extracts text from a PDF file, using OCR for image-based pages."""
    try:
        text = ""
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                else:
                    # if no text is extracted, use OCR
                    image = page.to_image()
                    text += pytesseract.image_to_string(image.original) + "\n"
        return text
    except Exception as e:
        raise ParsingError(f"Error parsing PDF: {e}")

def parse_docx(file_path: str) -> str:
    """Extracts text from a DOCX file."""
    try:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise ParsingError(f"Error parsing DOCX: {e}")

def parse_xlsx(file_path: str) -> str:
    """Extracts text from an XLSX file."""
    try:
        workbook = openpyxl.load_workbook(file_path)
        text = ""
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                text += "\t".join([str(cell.value) for cell in row if cell.value]) + "\n"
        return text
    except Exception as e:
        raise ParsingError(f"Error parsing XLSX: {e}")

def parse_pptx(file_path: str) -> str:
    """Extracts text from a PPTX file."""
    try:
        presentation = pptx.Presentation(file_path)
        text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise ParsingError(f"Error parsing PPTX: {e}")

def parse_text(file_path: str) -> str:
    """Extracts text from a plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        raise ParsingError(f"Error parsing text file: {e}")

def parse_csv(file_path: str) -> str:
    """Extracts text from a CSV file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            return "\n".join(["\t".join(row) for row in reader])
    except Exception as e:
        raise ParsingError(f"Error parsing CSV: {e}")

def parse_odt(file_path: str) -> str:
    """Extracts text from an ODT file."""
    try:
        doc = load(file_path)
        texts = []
        for element in doc.getElementsByType(text.P):
            texts.append(teletype.extractText(element))
        return "\n".join(texts)
    except Exception as e:
        raise ParsingError(f"Error parsing ODT: {e}")

def parse_rtf(file_path: str) -> str:
    """Extracts text from an RTF file."""
    try:
        with open(file_path, "r") as f:
            return rtf_to_text(f.read())
    except Exception as e:
        raise ParsingError(f"Error parsing RTF: {e}")

def parse_html(file_path: str) -> str:
    """Extracts text from an HTML file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
            return soup.get_text()
    except Exception as e:
        raise ParsingError(f"Error parsing HTML: {e}")

def parse_image(file_path: str) -> str:
    """Extracts text from an image file using OCR."""
    try:
        return pytesseract.image_to_string(file_path)
    except Exception as e:
        raise ParsingError(f"Error parsing image: {e}")

def parse_epub(file_path: str) -> str:
    """Extracts text from an EPUB file."""
    try:
        book = epub.read_epub(file_path)
        text = ""
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text += soup.get_text() + "\n"
        return text
    except Exception as e:
        raise ParsingError(f"Error parsing EPUB: {e}")

def parse_markdown(file_path: str) -> str:
    """Extracts text from a Markdown file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            md = MarkdownIt()
            html = md.render(f.read())
            soup = BeautifulSoup(html, "html.parser")
            return soup.get_text()
    except Exception as e:
        raise ParsingError(f"Error parsing Markdown: {e}")

# Main Processing Function
def process_document(file_path: str) -> (str, str):
    """
    Detects the MIME type of a file and calls the appropriate parser,
    with fallbacks for ambiguous text types.

    Args:
        file_path: The path to the file to process.

    Returns:
        A tuple containing the extracted text and the detected MIME type.
    """

    if os.path.getsize(file_path) == 0:
        return "", "inode/x-empty"

    mime_type = magic.from_file(file_path, mime=True)
    file_extension = os.path.splitext(file_path)[1].lower()

    parser_dispatcher_mime = {
        "application/pdf": parse_pdf,
        "image/jpeg": parse_image,
        "image/png": parse_image,
        "image/tiff": parse_image,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": parse_docx,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": parse_xlsx,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": parse_pptx,
        "text/plain": parse_text,
        "text/csv": parse_csv,
        "application/vnd.oasis.opendocument.text": parse_odt,
        "application/rtf": parse_rtf,
        "text/rtf": parse_rtf,
        "text/html": parse_html,
        "application/epub+zip": parse_epub,
        "text/markdown": parse_markdown,
    }

    parser = None

    # fallback for ambiguous text types
    if mime_type == "text/plain":
        if file_extension == ".csv":
            parser = parse_csv
            mime_type = "text/csv"
        elif file_extension == ".md":
            parser = parse_markdown
            mime_type = "text/markdown"
        else:
            parser = parse_text
    elif mime_type in parser_dispatcher_mime:
        parser = parser_dispatcher_mime[mime_type]

    if parser:
        raw_text = parser(file_path)
        cleaned_text = clean_text(raw_text)
        return cleaned_text, mime_type
    else:
        raise UnsupportedFileType(f"Unsupported file type: {mime_type}")
